# src/utils/pptx_generator.py

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from PySide6.QtGui import QFont
from PySide6.QtCore import QRectF

from app.models.song_model import Theme, Song
from.slide_layout_engine import split_lyrics_into_slides
from.text_formatter import PREFIXES_TO_HIGHLIGHT

def generate_presentation(songs: list, theme: Theme, output_path: str, overrides: dict):
    """
    Tạo một tệp PowerPoint.
    'overrides' giờ đây có thể chứa 'slides': list[str]
    """
    prs = Presentation()
    prs.slide_width = theme.slide_width
    prs.slide_height = theme.slide_height
    blank_slide_layout = prs.slide_layouts[6] 

    is_widescreen = theme.slide_width > 10000000
    slide_width_inches = Inches(13.333) if is_widescreen else Inches(10)
    slide_height_inches = Inches(7.5)

    for i, song in enumerate(songs):
        
        # Lấy override cho bài hát này
        song_override = overrides.get(song.id, {})
        
        lyric_size = song_override.get('lyric', theme.lyric_font_size)
        title_size = song_override.get('title', theme.title_font_size)
        
        lyrics_slides_for_export = []
        
        # --- LOGIC MỚI ĐỂ LẤY SLIDES ---
        if 'slides' in song_override:
            # ƯU TIÊN 1: Lấy các slide đã được người dùng chỉnh sửa
            lyrics_slides_for_export = song_override['slides']
            
            # (Chúng ta giả định slide đầu tiên vẫn là tựa đề,
            # và 'slides' chỉ chứa phần lời)
            # Chúng ta cần tách slide lời đầu tiên ra
            first_slide_lyric_chunk = lyrics_slides_for_export[0] if lyrics_slides_for_export else ""
            lyrics_slides_remaining = lyrics_slides_for_export[1:]
            
        else:
            # ƯU TIÊN 2: Dùng logic chia slide cũ (fallback)
            lyric_font = QFont(theme.lyric_font_name, lyric_size)
            slide_w_px, slide_h_px = (960, 540) if is_widescreen else (720, 540)
            
            # Box cho slide đầu (sau tựa đề)
            box_h_first = (slide_height_inches.emu - Inches(1).emu) / 914400.0 * 96
            temp_bounding_box = QRectF(0, 0, slide_w_px * 0.9, box_h_first * 0.9)
            temp_slides = split_lyrics_into_slides(song.lyrics, lyric_font, temp_bounding_box)
            
            first_slide_lyric_chunk = ""
            remaining_lyrics = song.lyrics
            
            if temp_slides:
                first_slide_lyric_chunk = temp_slides[0]
                remaining_lyrics = song.lyrics[len(first_slide_lyric_chunk):].lstrip()

            # Box cho các slide sau
            box_h_full = slide_height_inches.emu / 914400.0 * 96
            full_bounding_box = QRectF(0, 0, slide_w_px * 0.9, box_h_full * 0.9)
            lyrics_slides_remaining = split_lyrics_into_slides(remaining_lyrics, lyric_font, full_bounding_box)
            
        # --- HẾT LOGIC MỚI ---


        # --- 2. Tạo Slide Tựa đề (có lời) ---
        slide = prs.slides.add_slide(blank_slide_layout)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor.from_string(theme.bg_color[1:])

        # --- 2a. Thêm Textbox cho Tựa đề ---
        # (title_size đã được lấy từ override ở trên)
        txBox_title = slide.shapes.add_textbox(Inches(0), Inches(0), slide_width_inches, Inches(1))
        tf_title = txBox_title.text_frame
        p_title = tf_title.paragraphs[0]
        p_title.text = song.title # (Cũng có thể lấy từ override nếu bạn lưu cả tựa đề)
        p_title.alignment = PP_ALIGN.CENTER
        font_title = p_title.font
        font_title.name = theme.title_font_name
        font_title.size = Pt(title_size)
        font_title.color.rgb = RGBColor.from_string(theme.title_font_color[1:])
        font_title.bold = theme.title_font_bold
        font_title.italic = theme.title_font_italic
        font_title.underline = theme.title_font_underline

        # --- 2b. Thêm Textbox cho phần lời đầu tiên ---
        if first_slide_lyric_chunk:
            txBox_content = slide.shapes.add_textbox(Inches(0), Inches(1), slide_width_inches, slide_height_inches - Inches(1))
            tf_content = txBox_content.text_frame
            tf_content.word_wrap = True
            p_content = tf_content.paragraphs[0]
            _apply_lyric_formatting(p_content, first_slide_lyric_chunk, theme, lyric_size)

        # --- 3. Tạo các slide lời bài hát tiếp theo ---
        for slide_text in lyrics_slides_remaining:
            if not slide_text.strip(): continue
            
            slide = prs.slides.add_slide(blank_slide_layout)
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor.from_string(theme.bg_color[1:])

            txBox_full = slide.shapes.add_textbox(Inches(0), Inches(0), slide_width_inches, slide_height_inches)
            tf_full = txBox_full.text_frame
            tf_full.word_wrap = True
            p_full = tf_full.paragraphs[0]
            _apply_lyric_formatting(p_full, slide_text, theme, lyric_size)

        # --- 4. Slide chuyển tiếp ---
        if i < len(songs) - 1:
            slide = prs.slides.add_slide(blank_slide_layout)
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor.from_string("000000")

    prs.save(output_path)

# ... (Hàm _apply_lyric_formatting giữ nguyên)
def _apply_lyric_formatting(p: 'Paragraph', text: str, theme: Theme, lyric_size: int):
    """Hàm trợ giúp để áp dụng định dạng cho một đoạn văn bản lời bài hát."""
    # Đặt các thuộc tính chung cho cả đoạn văn
    if theme.lyric_alignment == "CENTER":
        p.alignment = PP_ALIGN.CENTER
    elif theme.lyric_alignment == "LEFT":
        p.alignment = PP_ALIGN.LEFT
    elif theme.lyric_alignment == "RIGHT":
        p.alignment = PP_ALIGN.RIGHT
    elif theme.lyric_alignment == "JUSTIFY":
        p.alignment = PP_ALIGN.JUSTIFY

    # Logic tô màu bằng cách thêm các "run"
    stripped_text = text.lstrip()
    found_prefix = None
    for prefix in PREFIXES_TO_HIGHLIGHT:
        if stripped_text.startswith(prefix):
            found_prefix = prefix
            break
    
    if found_prefix:
        run1 = p.add_run()
        run1.text = found_prefix
        font1 = run1.font
        font1.name = theme.title_font_name
        font1.color.rgb = RGBColor.from_string(theme.title_font_color[1:])
        font1.size = Pt(lyric_size)
        font1.bold = theme.lyric_font_bold
        font1.italic = theme.lyric_font_italic
        font1.underline = theme.lyric_font_underline
        
        run2 = p.add_run()
        run2.text = text[len(found_prefix):]
        font2 = run2.font
        font2.name = theme.lyric_font_name
        font2.size = Pt(lyric_size)
        font2.color.rgb = RGBColor.from_string(theme.lyric_font_color[1:])
        font2.bold = theme.lyric_font_bold
        font2.italic = theme.lyric_font_italic
        font2.underline = theme.lyric_font_underline
    else:
        run = p.add_run()
        run.text = text
        font = run.font
        font.name = theme.lyric_font_name
        font.size = Pt(lyric_size)
        font.color.rgb = RGBColor.from_string(theme.lyric_font_color[1:])
        font.bold = theme.lyric_font_bold
        font.italic = theme.lyric_font_italic
        font.underline = theme.lyric_font_underline